from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colours ──────────────────────────────────────────────────
BG        = RGBColor(0x08, 0x08, 0x0F)
PINK      = RGBColor(0xFF, 0x2D, 0x6E)
PURPLE    = RGBColor(0x7B, 0x61, 0xFF)
WHITE     = RGBColor(0xF0, 0xEE, 0xF8)
DIM       = RGBColor(0x88, 0x86, 0xA8)
GREEN     = RGBColor(0x00, 0xE5, 0xA0)
CARD      = RGBColor(0x12, 0x12, 0x22)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def txt(frame, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    p = frame.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Georgia" if bold and size >= 28 else "Calibri"
    return p

def add_box(slide, text, left, top, width, height, fontsize=11, bg=CARD, fg=DIM, bold=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = bg
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.color.rgb = fg
    run.font.bold = bold
    run.font.name = "Calibri"
    return box

def accent_bar(slide, top=0.18):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(top), Inches(13.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PINK
    bar.line.fill.background()

def so_what_box(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.7))
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x0A, 0x14)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "→  " + text
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = PINK
    run.font.name = "Calibri"

# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
accent_bar(slide, top=0.0)

title_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(1.4))
tf = title_box.text_frame
txt(tf, "We Can Move Faster", 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.3), Inches(10), Inches(0.7))
tf2 = sub_box.text_frame
txt(tf2, "AI-powered innovation cycles that compress months into days", 20, color=DIM, align=PP_ALIGN.CENTER)

bullets = [
    "A live demonstration built in a weekend — no engineering team required",
    "The same method is available to us now",
    "Our competitors are already moving this way",
]
bul_box = slide.shapes.add_textbox(Inches(3), Inches(4.2), Inches(7.5), Inches(2))
tf3 = bul_box.text_frame
tf3.word_wrap = True
for b in bullets:
    p = tf3.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "— " + b
    run.font.size = Pt(13)
    run.font.color.rgb = DIM
    run.font.name = "Calibri"
    p.space_before = Pt(4)

# Pink dot decoration
dot = slide.shapes.add_shape(9, Inches(6.4), Inches(1.6), Inches(0.12), Inches(0.12))
dot.fill.solid(); dot.fill.fore_color.rgb = PINK; dot.line.fill.background()

# ════════════════════════════════════════════════════════════
# SLIDE 2 — WHY NOW
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
accent_bar(slide)

head = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
txt(head.text_frame, "The Real Cost of Slow", 32, bold=True, color=WHITE)

bullets = [
    ("We have strong growth mandates and a full pipeline of ideas", WHITE),
    ("The constraint is how fast we can learn — validate, launch, measure, iterate", WHITE),
    ("Today: research takes weeks, tools take quarters, decisions lag behind the market", DIM),
    ("Every month of delay is a month a competitor learns what we don't", PINK),
]
bul_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.5))
bul_box.text_frame.word_wrap = True
for b, col in bullets:
    p = bul_box.text_frame.add_paragraph()
    run = p.add_run()
    run.text = "  ●  " + b
    run.font.size = Pt(15)
    run.font.color.rgb = col
    run.font.name = "Calibri"
    p.space_before = Pt(10)

so_what_box(slide, "The AI task force isn't protecting us — it's slowing us down while others accelerate.")

# ════════════════════════════════════════════════════════════
# SLIDE 3 — LEVELS
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
accent_bar(slide)

head = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
txt(head.text_frame, "Copilot vs. Factory — Two Very Different Positions", 28, bold=True, color=WHITE)

levels = [
    ("Level 1", "AI autocomplete  (GitHub Copilot)", DIM),
    ("Level 2", "AI writes functions", DIM),
    ("Level 3", "AI writes features", DIM),
    ("Level 4", "Human supervises agents  ← most advanced enterprise teams", WHITE),
    ("Level 5", "Autonomous software factory  ← this demo", PINK),
]

for i, (label, desc, col) in enumerate(levels):
    top = 1.3 + i * 0.88
    bg_col = RGBColor(0x20, 0x08, 0x14) if i == 4 else CARD
    lbox = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(1.5), Inches(0.65))
    lbox.fill.solid(); lbox.fill.fore_color.rgb = PINK if i == 4 else RGBColor(0x20,0x20,0x38)
    lp = lbox.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run(); lr.text = label
    lr.font.size = Pt(12); lr.font.bold = True
    lr.font.color.rgb = WHITE; lr.font.name = "Calibri"

    dbox = slide.shapes.add_textbox(Inches(2.4), Inches(top), Inches(9.5), Inches(0.65))
    dbox.fill.solid(); dbox.fill.fore_color.rgb = bg_col
    dp = dbox.text_frame.paragraphs[0]
    dp.alignment = PP_ALIGN.LEFT
    dr = dp.add_run(); dr.text = "  " + desc
    dr.font.size = Pt(13); dr.font.color.rgb = col; dr.font.name = "Calibri"
    dr.font.bold = (i == 4)

so_what_box(slide, "Most Fortune 100 orgs are at Level 2–3. Level 5 exists in production today.")

# ════════════════════════════════════════════════════════════
# SLIDE 4 — THE METHOD
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
accent_bar(slide)

head = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
txt(head.text_frame, "The Software Factory: Spec It, Run It, Learn From It", 28, bold=True, color=WHITE)

steps = [
    ("01", "Teams write what they need in plain English", "No IT ticket. No agency brief. Just a spec."),
    ("02", "An AI pipeline builds, tests, and delivers automatically", "Kilroy orchestrates AI agents through every stage."),
    ("03", "Humans govern at two points only", "Define the goal. Evaluate the result. AI does the rest."),
    ("04", "Failures loop back — the system self-corrects", "Git checkpoints at every stage. Fully auditable."),
]

for i, (num, title, sub) in enumerate(steps):
    left = 0.5 + i * 3.2
    nb = slide.shapes.add_textbox(Inches(left), Inches(1.3), Inches(2.9), Inches(0.5))
    nb.fill.solid(); nb.fill.fore_color.rgb = PINK
    np_ = nb.text_frame.paragraphs[0]
    np_.alignment = PP_ALIGN.CENTER
    nr = np_.add_run(); nr.text = num
    nr.font.size = Pt(18); nr.font.bold = True; nr.font.color.rgb = WHITE; nr.font.name = "Calibri"

    tb = slide.shapes.add_textbox(Inches(left), Inches(1.9), Inches(2.9), Inches(1.8))
    tb.fill.solid(); tb.fill.fore_color.rgb = CARD
    tf = tb.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run(); r1.text = title
    r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = WHITE; r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(11); r2.font.color.rgb = DIM; r2.font.name = "Calibri"
    p2.space_before = Pt(6)

add_box(slide,
    "Applied to CPG: concept testing tools, consumer survey platforms, launch trackers, "
    "campaign generators — built in days, not quarters.",
    0.5, 4.0, 12.3, 1.0, fontsize=13, bg=RGBColor(0x0C,0x08,0x1A), fg=WHITE)

so_what_box(slide, "Your innovation team becomes the builder. IT becomes the governance layer.")

# ════════════════════════════════════════════════════════════
# SLIDE 5 — THE DEMO
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
accent_bar(slide)

head = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
txt(head.text_frame, "What This Looks Like in Practice — Built in a Weekend", 28, bold=True, color=WHITE)

week = [
    ("MON", "Define hypothesis", PINK),
    ("TUE", "AI builds the tool", PURPLE),
    ("WED", "Consumer panel runs", GREEN),
    ("THU", "AI analyses & recommends", PINK),
    ("FRI", "Next concept is live", PURPLE),
]
for i, (day, act, col) in enumerate(week):
    left = 0.5 + i * 2.5
    db = slide.shapes.add_textbox(Inches(left), Inches(1.3), Inches(2.2), Inches(0.5))
    db.fill.solid(); db.fill.fore_color.rgb = col
    dp = db.text_frame.paragraphs[0]; dp.alignment = PP_ALIGN.CENTER
    dr = dp.add_run(); dr.text = day
    dr.font.size = Pt(13); dr.font.bold = True; dr.font.color.rgb = WHITE; dr.font.name = "Calibri"

    ab = slide.shapes.add_textbox(Inches(left), Inches(1.9), Inches(2.2), Inches(0.8))
    ab.fill.solid(); ab.fill.fore_color.rgb = CARD
    ap = ab.text_frame.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
    ar = ap.add_run(); ar.text = act
    ar.font.size = Pt(12); ar.font.color.rgb = WHITE; ar.font.name = "Calibri"

proof_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(12.3), Inches(2.6))
proof_box.fill.solid(); proof_box.fill.fore_color.rgb = RGBColor(0x10, 0x06, 0x1A)
tf = proof_box.text_frame; tf.word_wrap = True
lines = [
    ("PROOF OF CONCEPT: SNATCH", 12, PINK, True),
    ("A fully functional AI brand marketing platform — built in one weekend using this exact method.", 14, WHITE, False),
    ("Enter any brand's website URL → get brand DNA analysis + social media campaigns + AI-generated images.", 13, DIM, False),
    ("No engineering team. No agency. No quarters-long roadmap.", 13, GREEN, True),
]
for i, (line, sz, col, bold) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold; r.font.name = "Calibri"
    p.space_before = Pt(6)

so_what_box(slide, "A weekend. One person. A production-grade tool. That is the new baseline for our innovation org.")

# ════════════════════════════════════════════════════════════
# SLIDE 6 — THE ASK
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
accent_bar(slide)

head = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
txt(head.text_frame, "One Pilot. 30 Days. A New Innovation Pace.", 32, bold=True, color=WHITE)

# Left column
left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(5.8), Inches(4.5))
left_box.fill.solid(); left_box.fill.fore_color.rgb = CARD
ltf = left_box.text_frame; ltf.word_wrap = True
left_items = [
    ("THE SITUATION", 11, PINK, True),
    ("AI task force friction is costing us speed we cannot afford to lose.", 13, WHITE, False),
    ("Our competitors are compressing innovation cycles.", 13, DIM, False),
    ("We have the growth mandate. We need the velocity.", 13, WHITE, False),
]
for i, (line, sz, col, bold) in enumerate(left_items):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold; r.font.name = "Calibri"
    p.space_before = Pt(8)

# Right column
right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.8), Inches(4.5))
right_box.fill.solid(); right_box.fill.fore_color.rgb = RGBColor(0x14, 0x06, 0x20)
rtf = right_box.text_frame; rtf.word_wrap = True
right_items = [
    ("THE ASK", 11, PINK, True),
    ("1.  Designate one innovation workstream for a 30-day autonomous AI pilot.", 13, WHITE, False),
    ("2.  Remove task force friction for this pilot — measure the speed difference.", 13, WHITE, False),
    ("3.  Use results to define a new model for how our innovation org operates.", 13, WHITE, False),
]
for i, (line, sz, col, bold) in enumerate(right_items):
    p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold; r.font.name = "Calibri"
    p.space_before = Pt(10)

so_what_box(slide, "We don't need more resources. We need to remove the drag. This is how we get ahead of the growth pressure — and stay there.")

# ── Save ────────────────────────────────────────────────────
out = "/workspaces/claude-code-projects/workspace/projects/snatch-cio-pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
